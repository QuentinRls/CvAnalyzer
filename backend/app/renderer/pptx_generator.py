"""
Générateur de présentation PowerPoint avec template Devoteam

Ce module génère des présentations PowerPoint au format portrait (7.5" x 10") 
avec le design et la charte graphique Devoteam.

Caractéristiques principales:
- Format portrait optimisé pour l'impression et la présentation
- Utilisation des polices Montserrat (titres) et Montserrat Light (texte normal)
- Couleurs corporate Devoteam (rouge #F8485D, vert #009688)
- Layout en colonnes adaptatif selon le contenu
- Gestion automatique des bullet points PowerPoint (\n- format)
- Robustesse avec fallbacks en cas de données manquantes

Structure de la présentation:
1. Slide de titre : Identité + contact + langues + missions récentes
2. Slide de compétences : Techniques + fonctionnelles en colonnes
3. Slides d'expériences : Une slide détaillée par expérience professionnelle

Utilisation:
    from .pptx_generator import generate_devoteam_pptx
    
    pptx_file = generate_devoteam_pptx(dossier_competences)
    
Auteur: Générateur automatisé avec template Devoteam personnalisé
Version: 2.0 - Format portrait avec design specs exactes
"""
from io import BytesIO
from typing import List
import logging
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from ..schemas import DossierCompetences, ExperienceProfessionnelle

logger = logging.getLogger(__name__)

# Couleurs Devoteam
DEVOTEAM_RED = RGBColor(248, 72, 93)  # #F8485D
DEVOTEAM_GREEN = RGBColor(0, 150, 136)  # #009688
DARK_GRAY = RGBColor(51, 51, 51)  # #333333
LIGHT_GRAY = RGBColor(102, 102, 102)  # #666666
BACKGROUND_GRAY = RGBColor(245, 245, 245)  # #F5F5F5

# Chemin vers l'image Devoteam
DEVOTEAM_LOGO_PATH = os.path.join(os.path.dirname(__file__), "..", "assets", "devoteam.png")


class DevoteamPPTXGenerator:
    """
    Générateur de présentations PowerPoint avec template Devoteam
    
    Cette classe orchestre la création de présentations PowerPoint en format portrait
    selon les spécifications de design Devoteam. Elle gère:
    
    - Configuration du format de page (7.5" x 10", portrait)
    - Application de la charte graphique (polices, couleurs, logos)
    - Structure multi-slides avec navigation logique
    - Gestion des données manquantes avec fallbacks appropriés
    
    Fonctionnalités:
    - Auto-détection du contenu disponible (compétences, expériences, langues)
    - Mise en page adaptative selon la quantité de données
    - Formatage automatique des bullet points pour PowerPoint
    - Intégration du logo Devoteam (image PNG + fallback textuel)
    
    Standards de design:
    - Titres principaux: Montserrat gras, rouge Devoteam, 18-24pt
    - Sous-titres: Montserrat gras, gris foncé, 14-16pt  
    - Texte normal: Montserrat Light, noir/gris, 9pt
    - Bullet points: Format "\n- " pour compatibilité PowerPoint native
    """
    
    def __init__(self):
        """
        Initialise le générateur PowerPoint avec configuration portrait
        
        Configure automatiquement:
        - Dimensions de slide: 7.5" x 10" (portrait)
        - Layout de base: slides vierges pour contrôle total du positionnement
        - Template vide pour construction from scratch
        """
        # === CONFIGURATION DE BASE ===
        # Création d'une présentation vierge
        self.prs = Presentation()
        
        # === CONFIGURATION FORMAT PORTRAIT ===
        # Dimensions de slide: 7.5" largeur x 10" hauteur (format portrait)
        self.prs.slide_width = Inches(7.5)   # Largeur réduite pour format portrait
        self.prs.slide_height = Inches(10.0)  # Hauteur augmentée pour format portrait
        
        # Note: Cette configuration donne un ratio 3:4 optimisé pour l'impression
        # et la consultation sur écran en orientation portrait
        self.prs.slide_height = Inches(10)   # Format portrait (A4)
    
    # --- NOUVEAU : utilitaire pour tronquer au premier point ---
    def _truncate_at_first_dot(self, text: str) -> str:
        """
        Retourne la sous-chaîne jusqu'au premier point inclus (si présent),
        sinon retourne le texte en l'état (trim).
        """
        if not text:
            return ""
        idx = text.find('.')
        if idx == -1:
            return text.strip()
        return text[:idx + 1].strip()
    
    def generate_presentation(self, dossier: DossierCompetences) -> BytesIO:
        """
        Génère une présentation PowerPoint complète
        
        Args:
            dossier: Données structurées du CV
            
        Returns:
            BytesIO contenant le fichier PPTX
        """
        try:
            # Slide 1: Page de présentation
            self._create_title_slide(dossier)
            
            # Slide 2: Compétences
            self._create_skills_slide(dossier)
            
            # Slides 3+: Expériences détaillées
            if dossier.experiences_professionnelles:
                for exp in dossier.experiences_professionnelles:
                    self._create_experience_slide(exp)
            
            # Sauvegarder en BytesIO
            buffer = BytesIO()
            self.prs.save(buffer)
            buffer.seek(0)
            
            logger.info("Présentation PowerPoint générée avec succès")
            return buffer
            
        except Exception as e:
            logger.error(f"Erreur lors de la génération PowerPoint : {str(e)}")
            raise
    
    def _add_devoteam_logo(self, slide, x=Inches(0.0), y=Inches(0.0), size=Inches(0.5)):
        """
        Ajoute le logo Devoteam à une slide avec gestion robuste des erreurs
        
        Cette fonction tente d'abord d'utiliser l'image PNG du logo Devoteam.
        En cas d'échec (fichier manquant ou erreur de lecture), elle utilise
        un fallback textuel avec la même apparence visuelle.
        
        Args:
            slide: Slide PowerPoint où ajouter le logo
            x: Position horizontale (par défaut: coin gauche)
            y: Position verticale (par défaut: coin haut)
            size: Taille du logo (appliquée en largeur ET hauteur pour un format carré)
            
        Returns:
            Shape object (image ou textbox) contenant le logo
            
        Note: Le logo est forcé en format carré pour cohérence visuelle
        """
        try:
            # === TENTATIVE D'UTILISATION DE L'IMAGE PNG ===
            if os.path.exists(DEVOTEAM_LOGO_PATH):
                # Utiliser l'image réelle avec dimensions forcées en carré
                # Cela garantit un rendu cohérent même si l'image originale n'est pas carrée
                logo_shape = slide.shapes.add_picture(DEVOTEAM_LOGO_PATH, x, y, width=size*2.5, height=size)
                return logo_shape
            else:
                # === FALLBACK TEXTUEL ===
                logger.warning(f"Image Devoteam non trouvée à {DEVOTEAM_LOGO_PATH}, utilisation du texte")
                # Créer une textbox avec les mêmes dimensions carrées
                logo_box = slide.shapes.add_textbox(x, y, size, size)
                logo_frame = logo_box.text_frame
                logo_frame.text = "devoteam"
                logo_para = logo_frame.paragraphs[0]
                logo_para.font.size = Pt(20)  # Taille appropriée pour lisibilité
                logo_para.font.color.rgb = DEVOTEAM_RED  # Couleur rouge corporate
                logo_para.font.bold = True
                logo_para.font.name = "Montserrat"  # Police cohérente avec le reste
                return logo_box
                
        except Exception as e:
            # === FALLBACK D'URGENCE ===
            logger.error(f"Erreur lors de l'ajout du logo : {str(e)}")
            # En cas d'erreur critique, assurer qu'un logo textuel apparaît toujours
            logo_box = slide.shapes.add_textbox(x, y, size, size)
            logo_frame = logo_box.text_frame
            logo_frame.text = "devoteam"
            logo_para = logo_frame.paragraphs[0]
            logo_para.font.size = Pt(20)
            logo_para.font.color.rgb = DEVOTEAM_RED
            logo_para.font.bold = True
            return logo_box
    
    def _create_title_slide(self, dossier: DossierCompetences):
        """
        Crée la slide de titre avec présentation générale selon le design fourni
        
        Layout: 
        - Logo Devoteam en haut à gauche
        - Titre, sous-titre et nom juste en dessous du logo (côté gauche)
        - Description professionnelle en colonne gauche
        - Diplômes en colonne gauche (milieu)
        - Langues en colonne gauche (bas)
        - Mission Devoteam en colonne droite
        """
        slide_layout = self.prs.slide_layouts[6]  # Layout vide
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Fond blanc
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)  # Blanc
        
        # === AJOUT DU CERCLE JAUNE ÉNORME EN ARRIÈRE-PLAN ===
        # Cercle jaune énorme dans le coin haut gauche, derrière tout
        yellow_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-4.2), Inches(-8.3),  # Position top-left
            Inches(10), Inches(10)   # Dimensions: 2x2 inches pour un cercle énorme
        )
        yellow_circle.fill.solid()
        yellow_circle.fill.fore_color.rgb = RGBColor(239, 234, 220)  # Jaune
        yellow_circle.line.fill.background()  # Pas de bordure
        yellow_circle.shadow.inherit = False  # Remove any default shadow effect

        yellow_circle2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-1.5), Inches(8.7),  # Position top-left
            Inches(4), Inches(4)   # Dimensions: 2x2 inches pour un cercle énorme
        )
        yellow_circle2.fill.solid()
        yellow_circle2.fill.fore_color.rgb = RGBColor(239, 234, 220)  # Jaune
        yellow_circle2.line.fill.background()  # Pas de bordure
        yellow_circle2.shadow.inherit = False
        # === ZONE HAUT GAUCHE : LOGO + IDENTITÉ ===
        # Logo Devoteam en haut à gauche
        self._add_devoteam_logo(slide, Inches(0.0), Inches(0.0), Inches(0.5))
        
        # Titre principal - Position: juste sous le logo, aligné à gauche
        if dossier.entete and dossier.entete.intitule_poste:
            title_box = slide.shapes.add_textbox(Inches(0.0), Inches(0.5), Inches(6.0), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = dossier.entete.intitule_poste
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(14)  # Taille réduite pour être sous le logo
            title_para.font.color.rgb = DEVOTEAM_RED
            title_para.font.bold = True
            title_para.font.name = "Montserrat"  # Titres en Montserrat normal (gras)
        
        # Sous-titre années d'expérience - Position: sous le titre
        if dossier.entete and dossier.entete.annees_experience:
            subtitle_box = slide.shapes.add_textbox(Inches(0.0), Inches(0.79), Inches(4.0), Inches(0.4))
            subtitle_frame = subtitle_box.text_frame
            
            exp_text = dossier.entete.annees_experience
            if not any(word in exp_text.lower() for word in ['année', 'ans', 'expérience']):
                exp_text = f"{exp_text} ans d'expériences"
            
            subtitle_frame.text = exp_text
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(12)  # Taille réduite
            subtitle_para.font.color.rgb = DEVOTEAM_RED
            subtitle_para.font.bold = True
            subtitle_para.font.name = "Montserrat"  # Sous-titres en Montserrat normal
        
        # Nom complet - Position: sous le sous-titre
        if dossier.entete:
            nom_complet = f"{dossier.entete.prenom or ''} {dossier.entete.nom or ''}".strip().upper()
            if nom_complet:
                name_box = slide.shapes.add_textbox(Inches(0.0), Inches(1.1), Inches(5.0), Inches(0.5))
                name_frame = name_box.text_frame
                name_frame.text = nom_complet
                name_para = name_frame.paragraphs[0]
                name_para.font.size = Pt(14)  # Taille réduite
                name_para.font.color.rgb = DEVOTEAM_RED
                name_para.font.bold = True
                name_para.font.name = "Montserrat"  # Noms en Montserrat normal (gras)
        
        # === COLONNE GAUCHE : INFORMATIONS PERSONNELLES ===
        # Description professionnelle - seulement si disponible
        if dossier.entete and dossier.entete.resume_profil:
            desc_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(5), Inches(1))
            desc_frame = desc_box.text_frame
            desc_frame.text = dossier.entete.resume_profil
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(9)  # Texte normal en taille 9
            desc_para.font.color.rgb = DARK_GRAY
            desc_para.font.name = "Montserrat Light"  # Texte normal en Montserrat Light
        
        # Section Diplômes - seulement si disponible
        if dossier.diplomes:
            # Titre de section (style différent du texte normal)
            diplomes_title = slide.shapes.add_textbox(Inches(0.0), Inches(1.8), Inches(1.9), Inches(0.3))
            diplomes_title_frame = diplomes_title.text_frame
            diplomes_title_frame.text = "Diplômes"
            diplomes_title_para = diplomes_title_frame.paragraphs[0]
            diplomes_title_para.font.size = Pt(14)  # Titres de section plus grands
            diplomes_title_para.font.color.rgb = DEVOTEAM_RED
            diplomes_title_para.font.bold = True
            diplomes_title_para.font.name = "Montserrat"  # Titres en Montserrat normal
            diplomes_title_para.alignment = PP_ALIGN.CENTER

            # Contenu des diplômes
            y_pos = 2.2
            for diplome in dossier.diplomes:
                diplome_box = slide.shapes.add_textbox(Inches(0.0), Inches(y_pos), Inches(1.9), Inches(0.8))
                diplome_frame = diplome_box.text_frame
                # Centrage vertical du contenu dans la textbox
                diplome_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                 # Construction du texte du diplôme
                diplome_text = f"🎓 {diplome.intitule or ''}"
                if diplome.etablissement:
                    diplome_text += f"\n{diplome.etablissement}"
                if diplome.annee:
                    diplome_text += f"\n{diplome.annee}"
                
                diplome_frame.text = diplome_text
                # Application du style aux paragraphes du diplôme
                for para in diplome_frame.paragraphs:
                    # Alignement horizontal centré
                    para.alignment = PP_ALIGN.CENTER
                    para.font.size = Pt(9)  # Texte normal en taille 9
                    para.font.color.rgb = DARK_GRAY
                    para.font.name = "Montserrat Light"  # Texte normal en Montserrat Light
                
                y_pos += 1.0  # Espacement entre diplômes
        
        # Section Langues - seulement si disponible  
        if dossier.langues:
            # Titre de la section Langues
            lang_title = slide.shapes.add_textbox(Inches(0.45), Inches(y_pos), Inches(1.5), Inches(0.3))
            lang_title_frame = lang_title.text_frame
            lang_title_frame.text = "Langues"
            lang_title_para = lang_title_frame.paragraphs[0]
            lang_title_para.font.size = Pt(14)  # Titres de section
            lang_title_para.font.color.rgb = DEVOTEAM_RED
            lang_title_para.font.bold = True
            lang_title_para.font.name = "Montserrat"  # Titres en Montserrat normal
            
            # Afficher les langues en colonne (même x, y qui augmente)
            x_col = 0.7  # position horizontale fixe pour la colonne
            y_circle = y_pos + 0.4  # position verticale de départ (sous le titre)
            for langue in dossier.langues:
                # Cercle pour la langue
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x_col), Inches(y_circle),
                    Inches(0.5), Inches(0.5)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Fond blanc
                circle.line.color.rgb = DEVOTEAM_RED  # Bordure rouge Devoteam
                circle.line.width = Pt(2)
                
                # Texte du niveau dans le cercle
                niveau_text = "natif" if langue.niveau and "natif" in langue.niveau.lower() else "technique"
                circle_text = circle.text_frame
                circle_text.text = niveau_text
                circle_para = circle_text.paragraphs[0]
                circle_para.font.size = Pt(7)  # Texte petit dans le cercle
                circle_para.font.color.rgb = DEVOTEAM_RED
                circle_para.font.bold = True
                circle_para.alignment = PP_ALIGN.CENTER
                circle_para.font.name = "Montserrat"
                circle_text.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                # Label de la langue sous le cercle (centré)
                lang_label = slide.shapes.add_textbox(
                    Inches(x_col - 0.25), Inches(y_circle + 0.6),
                    Inches(1), Inches(0.25)
                )
                lang_label_frame = lang_label.text_frame
                lang_label_frame.text = langue.langue or ""
                lang_label_para = lang_label_frame.paragraphs[0]
                lang_label_para.font.size = Pt(9)
                lang_label_para.font.color.rgb = DARK_GRAY
                lang_label_para.alignment = PP_ALIGN.CENTER
                lang_label_para.font.name = "Montserrat Light"
                
                # Décalage vertical pour la langue suivante
                y_circle += 0.9

        # === COLONNE DROITE : EXPÉRIENCES CLÉS RÉCENTES ===
        # Section Mission Devoteam - seulement si disponible
        if dossier.experiences_cles_recentes:
            # Titre de la section Mission
            mission_title = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(5.0), Inches(0.3))
            mission_title_frame = mission_title.text_frame
            mission_title_frame.text = "Expériences clés récentes"
            mission_title_para = mission_title_frame.paragraphs[0]
            mission_title_para.font.size = Pt(14)  # Titre de section
            mission_title_para.font.color.rgb = DEVOTEAM_RED
            mission_title_para.font.bold = True
            mission_title_para.font.name = "Montserrat"  # Titres en Montserrat normal
            
            # Contenu des missions - utiliser les vraies données du CV
            mission_box = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(5.5), Inches(5.5))
            mission_frame = mission_box.text_frame
            mission_frame.clear()
            
            # Utiliser les expériences clés récentes du dossier
            for i, exp in enumerate(dossier.experiences_cles_recentes[:5]):  # Limiter à 3 expériences
                if i > 0:
                    # Saut de ligne entre expériences
                    para = mission_frame.add_paragraph()
                    para.text = ""
                
                # Titre de l'expérience
                title_text = ""
                if exp.intitule_poste:
                    title_text = exp.intitule_poste
                if exp.client:
                    title_text += f" - {exp.client}"
                
                if title_text:
                    if i == 0:
                        mission_frame.text = title_text
                        para = mission_frame.paragraphs[0]
                    else:
                        para = mission_frame.add_paragraph()
                        para.text = title_text
                    
                    # Style pour les titres d'expérience
                    para.font.size = Pt(9)  # Un peu plus grand que le texte normal
                    para.font.bold = True
                    para.font.color.rgb = DARK_GRAY
                    para.font.name = "Montserrat"  # Titres en Montserrat normal
                    para.font.underline = True

                # Durée de l'expérience
                if exp.duree:
                    para = mission_frame.add_paragraph()
                    para.text = f"Durée : {exp.duree}"
                    para.font.size = Pt(9)  # Texte plus petit pour la durée
                    para.font.color.rgb = DARK_GRAY
                    para.font.name = "Montserrat Light"  # Texte normal en Montserrat Light
                    para.font.italic = True

                # Description brève de l'expérience
                if exp.description_breve:
                    desc_text = exp.description_breve
                    # Ne conserver que la première phrase
                    para = mission_frame.add_paragraph()
                    para.text = self._truncate_at_first_dot(desc_text)
                    para.font.size = Pt(9)  # Texte normal en taille 9
                    para.font.color.rgb = DARK_GRAY
                    para.font.name = "Montserrat Light"  # Texte normal en Montserrat Light
                
                # Responsabilités (utiliser des vrais bullet points PowerPoint)
                if exp.responsabilites:
                    for resp in exp.responsabilites[:3]:  # Limiter à 3 responsabilités maximum
                        # Afficher uniquement la première phrase de chaque responsabilité
                        para = mission_frame.add_paragraph()
                        # run pour le bullet (gras)
                        bullet_run = para.add_run()
                        bullet_run.text = "     •"
                        bullet_run.font.size = Pt(9)
                        bullet_run.font.color.rgb = DARK_GRAY
                        bullet_run.font.name = "Montserrat"
                        bullet_run.font.bold = True

                        # run pour le texte de la responsabilité (normal)
                        text_run = para.add_run()
                        text_run.text = "       " + self._truncate_at_first_dot(resp)
                        text_run.font.size = Pt(9)
                        text_run.font.color.rgb = DARK_GRAY
                        text_run.font.name = "Montserrat Light"  # Police normale en Montserrat Light

    def _create_skills_slide(self, dossier: DossierCompetences):
        """
        Crée la slide des compétences selon le design fourni
        
        Layout:
        - En-tête: Logo + Titre + Nom (même que slide 1)
        - Section "Compétences techniques" avec sous-catégories:
          * Langages de programmation
          * Développement Web  
          * Back-End & API
          * etc.
        - Section "Compétences fonctionnelles" avec sous-catégories:
          * Gestion de projet & organisation
          * Analyse & résolution de problèmes
        
        Styles:
        - Titres: Montserrat normal, gras
        - Texte normal: Montserrat Light, taille 9pt
        """
        slide_layout = self.prs.slide_layouts[6]  # Layout vide
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Fond blanc
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)  # Blanc
        
        # === AJOUT DU CERCLE JAUNE ÉNORME EN ARRIÈRE-PLAN ===
        # Cercle jaune énorme dans le coin haut gauche, derrière tout
        yellow_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-4.2), Inches(-8.3),  # Position top-left
            Inches(10), Inches(10)   # Dimensions: 2x2 inches pour un cercle énorme
        )
        yellow_circle.fill.solid()
        yellow_circle.fill.fore_color.rgb = RGBColor(239, 234, 220)  # Jaune
        yellow_circle.line.fill.background() 
        yellow_circle.shadow.inherit = False  # Remove any default shadow effect
        yellow_circle2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-1.5), Inches(8.7),  # Position top-left
            Inches(4), Inches(4)   # Dimensions: 2x2 inches pour un cercle énorme
        )
        yellow_circle2.fill.solid()
        yellow_circle2.fill.fore_color.rgb = RGBColor(239, 234, 220)  # Jaune
        yellow_circle2.line.fill.background()  # Pas de bordure
        yellow_circle2.shadow.inherit = False
        # === EN-TÊTE : LOGO + IDENTITÉ (même que slide 1) ===
        # Logo Devoteam en haut à gauche (adjusted to match title slide for consistency)
        self._add_devoteam_logo(slide, Inches(0.0), Inches(0.0), Inches(0.5))
        
        # Titre principal (récupéré des données du CV) - matched to title slide style
        if dossier.entete and dossier.entete.intitule_poste:
            title_box = slide.shapes.add_textbox(Inches(0.0), Inches(0.5), Inches(6.0), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = dossier.entete.intitule_poste
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(14)  # Taille réduite pour être sous le logo
            title_para.font.color.rgb = DEVOTEAM_RED
            title_para.font.bold = True
            title_para.font.name = "Montserrat"  # Titres en Montserrat normal (gras)
        
        # Sous-titre années d'expérience - matched to title slide style
        if dossier.entete and dossier.entete.annees_experience:
            subtitle_box = slide.shapes.add_textbox(Inches(0.0), Inches(0.79), Inches(4.0), Inches(0.4))
            subtitle_frame = subtitle_box.text_frame
            
            exp_text = dossier.entete.annees_experience
            if not any(word in exp_text.lower() for word in ['année', 'ans', 'expérience']):
                exp_text = f"{exp_text} ans d'expériences"
            
            subtitle_frame.text = exp_text
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(12)  # Taille réduite
            subtitle_para.font.color.rgb = DEVOTEAM_RED
            subtitle_para.font.bold = True
            subtitle_para.font.name = "Montserrat"  # Sous-titres en Montserrat normal
        
        # Nom complet - matched to title slide style
        if dossier.entete:
            nom_complet = f"{dossier.entete.prenom or ''} {dossier.entete.nom or ''}".strip().upper()
            if nom_complet:
                name_box = slide.shapes.add_textbox(Inches(0.0), Inches(1.1), Inches(5.0), Inches(0.5))
                name_frame = name_box.text_frame
                name_frame.text = nom_complet
                name_para = name_frame.paragraphs[0]
                name_para.font.size = Pt(14)  # Taille réduite
                name_para.font.color.rgb = DEVOTEAM_RED
                name_para.font.bold = True
                name_para.font.name = "Montserrat"  # Noms en Montserrat normal (gras)
        
        # === SECTION COMPÉTENCES TECHNIQUES ===
        # Titre de section
        tech_title = slide.shapes.add_textbox(Inches(2), Inches(1.7), Inches(6.5), Inches(0.4))
        tech_title_frame = tech_title.text_frame
        tech_title_frame.text = "Compétences techniques :"
        tech_title_para = tech_title_frame.paragraphs[0]
        tech_title_para.font.size = Pt(12)  # Titre de section
        tech_title_para.font.color.rgb = DEVOTEAM_RED
        tech_title_para.font.bold = True
        tech_title_para.font.name = "Montserrat"  # Titres de section en Montserrat normal
        
        # Contenu des compétences techniques (seulement si disponible dans les données)
        if dossier.competences_techniques:
            comp_tech = dossier.competences_techniques
            y_pos = 2.0
            
            # Langages de programmation
            if comp_tech.language_framework:
                lang_title = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6.5), Inches(0.3))
                lang_title_frame = lang_title.text_frame
                lang_title_frame.text = "Langages de programmation :"
                lang_title_para = lang_title_frame.paragraphs[0]
                lang_title_para.font.size = Pt(9)
                lang_title_para.font.color.rgb = DARK_GRAY
                lang_title_para.font.bold = True
                lang_title_para.font.name = "Montserrat"
                
                y_pos += 0.2
                for i, lang in enumerate(comp_tech.language_framework):
                    lang_item = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos), Inches(6), Inches(0.25))
                    lang_item_frame = lang_item.text_frame
                    lang_item_frame.text = f"-        {lang}"
                    lang_item_para = lang_item_frame.paragraphs[0]
                    lang_item_para.font.size = Pt(9)
                    lang_item_para.font.color.rgb = DARK_GRAY
                    lang_item_para.font.name = "Montserrat"
                    y_pos += 0.18
                y_pos += 0.2
            
            # Développement Web
            web_skills = []
            if hasattr(comp_tech, 'tests') and comp_tech.tests:
                web_skills.extend(comp_tech.tests)
            if hasattr(comp_tech, 'outils') and comp_tech.outils:
                web_skills.extend(comp_tech.outils)
            y_pos-= 0.1
            if web_skills:
                web_title = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6.5), Inches(0.3))
                web_title_frame = web_title.text_frame
                web_title_frame.text = "Développement Web :"
                web_title_para = web_title_frame.paragraphs[0]
                web_title_para.font.size = Pt(9)
                web_title_para.font.color.rgb = DARK_GRAY
                web_title_para.font.bold = True
                web_title_para.font.name = "Montserrat"
                
                y_pos += 0.2
                for skill in web_skills[:4]:  # Limiter à 4 éléments
                    skill_item = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos), Inches(6), Inches(0.25))
                    skill_item_frame = skill_item.text_frame
                    skill_item_frame.text = f"-        {skill}"
                    skill_item_para = skill_item_frame.paragraphs[0]
                    skill_item_para.font.size = Pt(9)
                    skill_item_para.font.color.rgb = DARK_GRAY
                    skill_item_para.font.name = "Montserrat"
                    y_pos += 0.18
                y_pos += 0.1
            
            # Back-End & API
            if comp_tech.base_de_donnees_big_data or comp_tech.ci_cd:
                backend_title = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6.5), Inches(0.3))
                backend_title_frame = backend_title.text_frame
                backend_title_frame.text = "Back-End & API :"
                backend_title_para = backend_title_frame.paragraphs[0]
                backend_title_para.font.size = Pt(9)
                backend_title_para.font.color.rgb = DARK_GRAY
                backend_title_para.font.bold = True
                backend_title_para.font.name = "Montserrat"
                
                y_pos += 0.2
                all_backend = []
                if comp_tech.base_de_donnees_big_data:
                    all_backend.extend(comp_tech.base_de_donnees_big_data)
                if comp_tech.ci_cd:
                    all_backend.extend(comp_tech.ci_cd)
                
                for skill in all_backend[:3]:
                    skill_item = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos), Inches(6), Inches(0.25))
                    skill_item_frame = skill_item.text_frame
                    skill_item_frame.text = f"- {skill}"
                    skill_item_para = skill_item_frame.paragraphs[0]
                    skill_item_para.font.size = Pt(9)
                    skill_item_para.font.color.rgb = DARK_GRAY
                    skill_item_para.font.name = "Montserrat"
                    y_pos += 0.18

        # Section Compétences fonctionnelles (seulement si disponible)
            y_pos -= 0.3
        if dossier.competences_fonctionnelles and y_pos < 8.5:
            comp_func = dossier.competences_fonctionnelles
            
            y_pos += 0.4
            func_title = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6.5), Inches(0.4))
            func_title_frame = func_title.text_frame
            func_title_frame.text = "Compétences fonctionnelles :"
            func_title_para = func_title_frame.paragraphs[0]
            func_title_para.font.size = Pt(12)
            func_title_para.font.color.rgb = DEVOTEAM_RED
            func_title_para.font.bold = True
            func_title_para.font.name = "Montserrat"
            
            y_pos += 0.3
            
            # Gestion de projet & organisation (seulement si disponible)
            if comp_func.gestion_de_projet:
                proj_title = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6.5), Inches(0.3))
                proj_title_frame = proj_title.text_frame
                proj_title_frame.text = "Gestion de projet & organisation :"
                proj_title_para = proj_title_frame.paragraphs[0]
                proj_title_para.font.size = Pt(9)
                proj_title_para.font.color.rgb = DARK_GRAY
                proj_title_para.font.bold = True
                proj_title_para.font.name = "Montserrat"
                
                y_pos += 0.2
                for skill in comp_func.gestion_de_projet:
                    skill_item = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos), Inches(6), Inches(0.25))
                    skill_item_frame = skill_item.text_frame
                    skill_item_frame.text = f"- {skill}"
                    skill_item_para = skill_item_frame.paragraphs[0]
                    skill_item_para.font.size = Pt(9)
                    skill_item_para.font.color.rgb = DARK_GRAY
                    skill_item_para.font.name = "Montserrat"
                    y_pos += 0.18
                
                y_pos += 0.2
            
            # Analyse & résolution de problèmes (seulement si disponible)
            if hasattr(comp_func, 'analyse_problemes') or comp_func.methodologie_scrum:
                analysis_title = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6.5), Inches(0.3))
                analysis_title_frame = analysis_title.text_frame
                analysis_title_frame.text = "Analyse & résolution de problèmes :"
                analysis_title_para = analysis_title_frame.paragraphs[0]
                analysis_title_para.font.size = Pt(12)
                analysis_title_para.font.color.rgb = DARK_GRAY
                analysis_title_para.font.bold = True
                analysis_title_para.font.name = "Montserrat"
                
                y_pos += 0.2
                # Ajouter les méthodologies si disponibles
                if comp_func.methodologie_scrum:
                    for method in comp_func.methodologie_scrum:
                        skill_item = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos), Inches(6), Inches(0.25))
                        skill_item_frame = skill_item.text_frame
                        skill_item_frame.text = f"- {method}"
                        skill_item_para = skill_item_frame.paragraphs[0]
                        skill_item_para.font.size = Pt(9)
                        skill_item_para.font.color.rgb = DARK_GRAY
                        skill_item_para.font.name = "Montserrat"
                        y_pos += 0.25
    
    def _create_experience_slide(self, exp: ExperienceProfessionnelle):
        """
        Crée une slide détaillée pour une expérience professionnelle
        
        Layout:
        - En-tête: Logo Devoteam + Titre "Expériences professionnelles récentes"
        - Section principale:
          * Logo/nom de l'entreprise cliente (placeholder)
          * Titre du poste et informations (période, lieu)
          * Contexte de la mission
          * Responsabilités principales (en bullet points)
          * Technologies utilisées
        
        Styles:
        - Titre principal: Montserrat, rouge Devoteam, gras, 24pt
        - Sous-titres: Montserrat, gras, 14-16pt
        - Texte normal: Montserrat Light, 9pt, noir/gris foncé
        - Bullet points: Format "\n- " pour compatibilité PowerPoint
        """
        slide_layout = self.prs.slide_layouts[6]  # Layout vide pour contrôle total
        slide = self.prs.slides.add_slide(slide_layout)
        
        # === CONFIGURATION DE BASE ===
        # Fond blanc uniforme
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)  # Blanc pur
        
        # === AJOUT DU CERCLE JAUNE ÉNORME EN ARRIÈRE-PLAN ===
        # Cercle jaune énorme dans le coin haut gauche, derrière tout
        yellow_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-4.2), Inches(-8.3),  # Position top-left
            Inches(10), Inches(10)   # Dimensions: 2x2 inches pour un cercle énorme
        )
        yellow_circle.fill.solid()
        yellow_circle.fill.fore_color.rgb = RGBColor(239, 234, 220)  # Jaune
        yellow_circle.line.fill.background() 
        yellow_circle.shadow.inherit = False  # Remove any default shadow effect
        
        # === CONFIGURATION DE BASE ===
        # Fond blanc uniforme
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)  # Blanc pur
        
        # === EN-TÊTE ===
        # Logo Devoteam repositionné en haut à gauche
        self._add_devoteam_logo(slide, Inches(0.2), Inches(0.1), Inches(0.5))
        
        # Titre principal de la slide
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(6.5), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Expériences professionnelles récentes."
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)  # Taille cohérente avec les autres slides
        title_para.font.color.rgb = DEVOTEAM_RED
        title_para.font.bold = True
        title_para.font.name = "Montserrat"  # Titres en Montserrat normal
        
        # === SECTION ENTREPRISE ===
        # Logo/nom de l'entreprise (placeholder visuel)
        company_name = exp.client or "CLIENT"
        logo_text = company_name.upper() if len(company_name) > 6 else company_name.upper()
        company_logo = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(1.5), Inches(1))
        company_logo_frame = company_logo.text_frame
        company_logo_frame.text = logo_text
        company_logo_para = company_logo_frame.paragraphs[0]
        company_logo_para.font.size = Pt(14)
        company_logo_para.font.color.rgb = LIGHT_GRAY  # Gris clair pour effet placeholder
        company_logo_para.font.bold = True
        company_logo_para.alignment = PP_ALIGN.CENTER
        company_logo_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Centrage vertical
        
        # Informations de l'expérience
        info_box = slide.shapes.add_textbox(Inches(2.2), Inches(1.5), Inches(4.8), Inches(1))
        info_frame = info_box.text_frame
        
        info_text = ""
        if exp.client:
            info_text += f"Client\n"
        if exp.intitule_poste:
            info_text += f"{exp.intitule_poste}\n"
        if exp.date_debut or exp.date_fin:
            duration = f"{exp.date_debut or ''} à {exp.date_fin or ''}".strip()
            info_text += f"{duration}"
        
        info_frame.text = info_text
        for i, para in enumerate(info_frame.paragraphs):
            para.font.name = "Montserrat"
            if i == 0:  # "Client"
                para.font.size = Pt(16)
                para.font.color.rgb = DARK_GRAY
                para.font.bold = True
            elif i == 1:  # Poste
                para.font.size = Pt(14)
                para.font.color.rgb = DEVOTEAM_RED
                para.font.bold = True
            else:  # Durée
                para.font.size = Pt(12)
                para.font.color.rgb = LIGHT_GRAY
        
        # === SECTION CONTEXTE ===
        # Affichage du contexte seulement si disponible dans les données
        if exp.contexte:
            # Titre de la section contexte
            context_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(6.5), Inches(0.3))
            context_title_frame = context_title.text_frame
            context_title_frame.text = "Contexte."
            context_title_para = context_title_frame.paragraphs[0]
            context_title_para.font.size = Pt(16)  # Titre de section
            context_title_para.font.color.rgb = DEVOTEAM_RED
            context_title_para.font.bold = True
            context_title_para.font.name = "Montserrat"  # Titres en Montserrat normal
            
            # Contenu du contexte
            context_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(6.5), Inches(1.2))
            context_frame = context_box.text_frame
            # Ne conserver que la première phrase du contexte
            context_frame.text = self._truncate_at_first_dot(exp.contexte)
            context_para = context_frame.paragraphs[0]
            context_para.font.size = Pt(12)  # Texte de contexte légèrement plus grand
            context_para.font.color.rgb = DARK_GRAY
            context_para.font.name = "Montserrat Light"  # Texte normal en Montserrat Light
        
        # === SECTION RESPONSABILITÉS ===
        # Affichage des responsabilités seulement si disponibles
        if exp.responsabilites:
            # Titre de la section responsabilités
            resp_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(6.5), Inches(0.3))
            resp_title_frame = resp_title.text_frame
            resp_title_frame.text = "Responsabilités."
            resp_title_para = resp_title_frame.paragraphs[0]
            resp_title_para.font.size = Pt(16)
            resp_title_para.font.color.rgb = DEVOTEAM_RED
            resp_title_para.font.bold = True
            resp_title_para.font.name = "Montserrat"  # Titres en Montserrat normal
            
            # Liste des responsabilités avec bullet points PowerPoint
            # Conserver uniquement la première phrase de chaque responsabilité
            resp_content = "\n- ".join(self._truncate_at_first_dot(r) for r in exp.responsabilites)
            resp_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(6.5), Inches(2.0))
            resp_frame = resp_box.text_frame
            resp_frame.text = f"- {resp_content}"  # Commencer par un tiret
            
            # Application du style à chaque bullet point
            for para in resp_frame.paragraphs:
                para.font.size = Pt(9)  # Taille 9 pour le texte normal (cohérent avec design)
                para.font.color.rgb = DARK_GRAY
                para.font.name = "Montserrat Light"  # Texte normal en Montserrat Light
        
        # === SECTION LIVRABLES ===
        # Affichage des livrables seulement si disponibles
        if exp.livrables:
            # Titre de la section livrables
            deliv_title = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(6.5), Inches(0.3))
            deliv_title_frame = deliv_title.text_frame
            deliv_title_frame.text = "Livrables."
            deliv_title_para = deliv_title_frame.paragraphs[0]
            deliv_title_para.font.size = Pt(16)  # Titre de section
            deliv_title_para.font.color.rgb = DEVOTEAM_RED
            deliv_title_para.font.bold = True
            deliv_title_para.font.name = "Montserrat"  # Titres en Montserrat normal
            
            # Liste des livrables avec bullet points PowerPoint
            # Conserver uniquement la première phrase de chaque livrable
            deliv_content = "\n- ".join(self._truncate_at_first_dot(l) for l in exp.livrables)
            deliv_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.6), Inches(6.5), Inches(2.0))
            deliv_frame = deliv_box.text_frame
            deliv_frame.text = f"- {deliv_content}"  # Commencer par un tiret
            
            # Application du style à chaque bullet point
            for para in deliv_frame.paragraphs:
                para.font.size = Pt(9)  # Taille 9 pour le texte normal
                para.font.color.rgb = DARK_GRAY
                para.font.name = "Montserrat Light"  # Texte normal en Montserrat Light


def generate_devoteam_pptx(dossier: DossierCompetences) -> BytesIO:
    """
    Génère une présentation PowerPoint avec le template Devoteam
    
    Args:
        dossier: Données structurées du CV
        
    Returns:
        BytesIO contenant le fichier PPTX
    """
    try:
        generator = DevoteamPPTXGenerator()
        return generator.generate_presentation(dossier)
        
    except Exception as e:
        logger.error(f"Erreur lors de la génération PowerPoint : {str(e)}")
        raise
